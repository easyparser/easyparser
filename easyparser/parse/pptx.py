import io
import logging
import re
import shutil
import subprocess
import tempfile
from operator import attrgetter
from pathlib import Path
from typing import Any

from easyparser.base import BaseOperation, Chunk, ChunkGroup, CType

logger = logging.getLogger(__name__)


def pptx_to_pdf(file_path: str, temp_dir: str) -> str:
    """Convert pptx to pdf using LibreOffice"""
    if not temp_dir:
        temp_dir = tempfile.gettempdir()

    file_name = Path(file_path).name
    output_dir = Path(temp_dir) / "output"
    soffice_dir = Path(temp_dir) / "soffice"

    print(f"Converting {file_path} to PDF at {output_dir}...")
    subprocess.run(
        [
            "soffice",
            f"-env:UserInstallation=file://{soffice_dir}",
            "--headless",
            "--convert-to",
            "pdf",
            file_path,
            "--outdir",
            str(output_dir),
        ],
        check=True,
    )

    return str(output_dir / file_name.replace(".pptx", ".pdf"))


def weighted_std(values, weights):
    """Return weighted standard deviation of the values"""
    import math

    import numpy as np

    if sum(weights) == 0:
        return 0.0

    average = np.average(values, weights=weights)
    variance = np.average((values - average) ** 2, weights=weights)

    return math.sqrt(variance)


def use_vlm_or_not(shapes, slide_width, slide_height) -> bool:
    """Decide whether should use VLM to generate a description of the slide.

    Based on following heuristics:
        - Check if the text locations spread out too much (hinting text location
        plays an important role to understand the slide)
    """
    x1s, weights = [], []
    for _i, shape in enumerate(shapes):
        if not hasattr(shape, "text"):
            continue
        x1s.append(shape.left / slide_width)
        weights.append(len(shape.text))

    return weighted_std(x1s, weights) > 0.15


def llm_caption(pil_image, gemini_api_key, gemini_model, prompt) -> str:
    """Generate a caption for the image using Gemini LLM."""
    from google import genai

    client = genai.Client(api_key=gemini_api_key)
    response = client.models.generate_content(
        model=gemini_model or "gemini-2.0-flash",
        contents=[pil_image, prompt],
    )
    return response.text or ""


def parse_image(shape, parent_chunk, slide_width, slide_height, **kwargs) -> Chunk:
    """Parse image, graph, chart, or other visual element. Combine provided alt
    text with LLM captioning.
    """
    from PIL import Image

    llm_description = ""
    alt_text = ""

    chunk = Chunk(
        mimetype=shape.image.content_type,
        content=shape.image.blob,
        ctype=CType.Figure,
        origin=parent_chunk.origin,
    )

    if llm_key := kwargs.get("gemini_api_key"):
        pil_img = Image.open(io.BytesIO(shape.image.blob))
        try:
            llm_description = llm_caption(
                pil_img,
                llm_key,
                kwargs.get("gemini_model"),
                "Describe this image. If there are any text inside this image, "
                "transcribe the text as closest as possible.",
            )
        except Exception as e:
            logger.error(f"Error in LLM captioning: {e}")
            llm_description = ""

    try:
        # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
    except Exception:
        pass

    alt_text = "\n".join([llm_description, alt_text]) or shape.name
    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
    alt_text = re.sub(r"\s+", " ", alt_text).strip()
    chunk.text = alt_text

    return chunk


def parse_table(shape, parent_chunk, slide_width, slide_height, **kwargs) -> Chunk:
    """Parse the table shape and convert to Markdown."""
    chunk = Chunk(
        mimetype="text/plain",
        ctype=CType.Figure,
        origin=parent_chunk.origin,
    )

    rows = []
    new_line = re.compile(r"\r?\n")
    for idx, row in enumerate(shape.table.rows, start=0):
        row_str = ""
        for cell in row.cells:
            cell_text = cell.text.replace("\\|", "[[ESCAPED_PIPE]]")
            cell_text = cell_text.replace("|", "\\|")
            cell_text = cell_text.replace("[[ESCAPED_PIPE]]", "\\|")
            cell_text = new_line.sub(" ", cell_text)
            row_str += f"| {cell_text} "
        row_str += "|"
        rows.append(row_str)
        if idx == 0:
            row_str += "\n" + "|".join(["---"] * len(row.cells)) + "|"

    chunk.text = "\n".join(rows)

    return chunk


def parse_chart(shape, parent_chunk, slide_width, slide_height, **kwargs) -> Chunk:
    """Parse the chart shape and convert to Markdown."""
    chunk = Chunk(
        mimetype="text/plain",
        ctype=CType.Para,
        origin=parent_chunk.origin,
    )
    chart = shape.chart
    md = "\n\n### Chart"
    if chart.has_title:
        md += f": {chart.chart_title.text_frame.text}"
    md += "\n\n"
    data = []
    category_names = [c.label for c in chart.plots[0].categories]
    series_names = [s.name for s in chart.series]
    data.append(["Category"] + series_names)

    for idx, category in enumerate(category_names):
        row = [category]
        for series in chart.series:
            row.append(series.values[idx])
        data.append(row)

    markdown_table = []
    for row in data:
        markdown_table.append("| " + " | ".join(map(str, row)) + " |")
    header = markdown_table[0]
    separator = "|" + "|".join(["---"] * len(data[0])) + "|"
    chunk.content = md + "\n".join([header, separator] + markdown_table[1:])
    return chunk


def parse_text(shape, parent_chunk, slide_width, slide_height, **kwargs) -> Chunk:
    """Parse the text shape"""
    chunk = Chunk(
        mimetype="text/plain",
        ctype=CType.Para,
        origin=parent_chunk.origin,
    )

    shape_text = ""
    for paragraph in shape.text_frame.paragraphs:
        level = paragraph.level or 0
        is_list = bool(paragraph._p.xpath("./a:pPr/a:buChar"))
        if is_list:
            text = " " * level + "- " + paragraph.text
        else:
            text = " " * level + paragraph.text

        if text:
            shape_text += text + "\n"

    chunk.content = shape_text.strip()
    return chunk


def parse_shape(
    shape, parent_chunk, slide_width, slide_height, **kwargs
) -> Chunk | None:
    import pptx.enum.shapes as ps

    if shape.shape_type == ps.MSO_SHAPE_TYPE.PICTURE:
        return parse_image(shape, parent_chunk, slide_width, slide_height, **kwargs)
    elif shape.shape_type == ps.MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image"):
        return parse_image(shape, parent_chunk, slide_width, slide_height, **kwargs)
    elif shape.shape_type == ps.MSO_SHAPE_TYPE.TABLE:
        return parse_table(shape, parent_chunk, slide_width, slide_height, **kwargs)
    elif shape.has_chart:
        return parse_chart(shape, parent_chunk, slide_width, slide_height, **kwargs)
    elif shape.has_text_frame:
        return parse_text(shape, parent_chunk, slide_width, slide_height, **kwargs)
    elif shape.shape_type == ps.MSO_SHAPE_TYPE.GROUP:
        sorted_shapes = sorted(shape.shapes, key=attrgetter("top", "left"))
        chunk = Chunk(
            mimetype="text/plain",
            ctype=CType.Div,
            origin=parent_chunk.origin,
        )
        childs = []
        for subshape in sorted_shapes:
            child_chunk = parse_shape(
                shape=subshape,
                parent_chunk=chunk,
                slide_width=slide_width,
                slide_height=slide_height,
                **kwargs,
            )
            if child_chunk:
                childs.append(child_chunk)
                child_chunk.parent = chunk
        for idx, child in enumerate(childs):
            if idx == 0:
                chunk.child = child
                continue
            child.prev = childs[idx - 1]
            childs[idx - 1].next = child

        if childs:
            return chunk

    return None


class PptxParser(BaseOperation):

    @classmethod
    def run(
        cls,
        chunk: Chunk | ChunkGroup,
        gemini_api_key: str | None = None,
        gemini_model: str | None = None,
        **kwargs: Any,
    ) -> ChunkGroup:
        import pptx
        import pypdfium2 as pdfium

        # Resolve chunk
        if isinstance(chunk, Chunk):
            chunk = ChunkGroup(chunks=[chunk])

        output = ChunkGroup()
        for mc in chunk:
            pres = pptx.Presentation(mc.origin.location)
            pdf_converted_path, temp_dir, pdf = "", "", None
            prev_slide = None
            for page_num, slide in enumerate(pres.slides, start=1):
                slide_children = []
                shapes = sorted(slide.shapes, key=attrgetter("top", "left"))

                use_vlm = False
                if gemini_api_key and use_vlm_or_not(
                    shapes, pres.slide_width, pres.slide_height
                ):
                    print(f"Using VLM for slide {page_num}")
                    use_vlm = True
                    if pdf is None:
                        temp_dir = tempfile.mkdtemp()
                        pdf_converted_path = pptx_to_pdf(mc.origin.location, temp_dir)
                        pdf = pdfium.PdfDocument(pdf_converted_path)

                    pil_image = pdf[page_num - 1].render().to_pil()
                    caption = llm_caption(
                        pil_image,
                        gemini_api_key,
                        gemini_model,
                        "Provide detailed description of this slide.",
                    )
                    slide_chunk = Chunk(
                        mimetype=CType.Div,
                        ctype=CType.Para,
                        text=f"**Slide {page_num}**\n{caption}",
                        origin=mc.origin,
                        parent=mc,
                        metadata={
                            "page_num": page_num,
                        },
                    )
                else:
                    slide_chunk = Chunk(
                        mimetype=CType.Div,
                        ctype=CType.Div,
                        origin=mc.origin,
                        parent=mc,
                        metadata={
                            "page_num": page_num,
                        },
                    )

                    for shape in shapes:
                        child = parse_shape(
                            shape=shape,
                            parent_chunk=slide_chunk,
                            slide_width=pres.slide_width,
                            slide_height=pres.slide_height,
                            gemini_api_key=gemini_api_key,
                            gemini_model=gemini_model,
                        )
                        if child is not None:
                            slide_children.append(child)
                            child.parent = slide_chunk

                if slide.has_notes_slide:
                    text_frame = ""
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame is not None:
                        text_frame = notes_frame.text

                    if text_frame:
                        child = Chunk(
                            mimetype="text/plain",
                            ctype=CType.Para,
                            origin=slide_chunk.origin,
                            content=f"**Slide note**: {text_frame}",
                        )
                        slide_children.append(child)
                        child.parent = slide_chunk
                        if use_vlm:
                            slide_chunk.text += f"\n{child.content}"

                for idx, child in enumerate(slide_children):
                    if idx == 0:
                        slide_chunk.child = child
                        continue
                    child.prev = slide_children[idx - 1]
                    slide_children[idx - 1].next = child

                if not use_vlm:
                    slide_chunk.text = f"**Slide {page_num}**{slide_chunk.render()}"

                if prev_slide is None:
                    mc.child = slide_chunk
                else:
                    slide_chunk.prev = prev_slide
                    prev_slide.next = slide_chunk

                prev_slide = slide_chunk

            if temp_dir:
                shutil.rmtree(temp_dir)
            output.append(mc)

        return output

    @classmethod
    def py_dependency(cls) -> list[str]:
        """Return the list of Python dependencies required by this converter."""
        return ["python-pptx", "google-genai", "Pillow", "pypdfium2", "numpy"]
